"""
Generate Demo Slide Deck (PPTX) for COMP 333 — 5-minute live demo.
6 slides with per-slide time cues.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = "Demo_Slides_TeamD.pptx"

# ── Palette ───────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x3A, 0x5C)
MID_BLUE   = RGBColor(0x2E, 0x6D, 0xA4)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
ACCENT     = RGBColor(0xE8, 0xA0, 0x20)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GREY  = RGBColor(0x44, 0x44, 0x44)
GREEN      = RGBColor(0x27, 0xAE, 0x60)
RED        = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────
def rect(slide, x, y, w, h, rgb):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = rgb
    sh.line.fill.background()
    return sh

def tb(slide, text, x, y, w, h,
       size=14, bold=False, italic=False,
       color=WHITE, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = box.text_frame; tf.word_wrap = True
    p   = tf.paragraphs[0]; p.alignment = align
    r   = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return box

def chrome(slide, timer):
    """Top bar + bottom bar + timer badge."""
    rect(slide, 0, 0, 13.33, 0.62, DARK_BLUE)
    tb(slide, "COMP 333 — End-to-End Data Analytics Pipeline · NYC Yellow Taxi",
       0.22, 0.08, 10, 0.48, size=13, color=WHITE)
    # Timer badge
    rect(slide, 11.55, 0.08, 1.55, 0.46, ACCENT)
    tb(slide, timer, 11.55, 0.08, 1.55, 0.46,
       size=13, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    # Footer
    rect(slide, 0, 7.1, 13.33, 0.4, DARK_BLUE)
    tb(slide, "Ronnie Chan · Patrice Gallant · Nesrine Larbi   |   Concordia University · April 2026",
       0.2, 7.13, 12.9, 0.33, size=10, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

def section_title(slide, text):
    tb(slide, text, 0.45, 0.76, 12.5, 0.62,
       size=24, bold=True, color=DARK_BLUE)
    rect(slide, 0.45, 1.42, 12.43, 0.04, ACCENT)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover  (talk for ≤ 20 s while it's showing)
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, DARK_BLUE)
rect(s, 8.8, 0, 4.53, 7.5, MID_BLUE)       # right panel
rect(s, 0,   3.25, 13.33, 0.06, ACCENT)    # accent stripe

tb(s, "COMP 333", 0.6, 1.05, 8, 0.55, size=16, color=LIGHT_BLUE)
tb(s, "End-to-End Data\nAnalytics Pipeline",
   0.6, 1.6, 8.5, 1.55, size=40, bold=True, color=WHITE)
tb(s, "NYC Yellow Taxi Trip Records — Jun & Jul 2025",
   0.6, 3.15, 8.5, 0.52, size=18, color=LIGHT_BLUE)

for i, m in enumerate(["Ronnie Chan (27206003)",
                         "Patrice Gallant (40301020)",
                         "Nesrine Larbi (40079009)"]):
    tb(s, m, 0.6, 3.95 + i*0.52, 5.5, 0.46, size=15, color=WHITE)

tb(s, "Concordia University · April 2026",
   0.6, 5.82, 5.5, 0.4, size=13, color=LIGHT_BLUE)

# Right panel timing guide
tb(s, "5-min Demo Plan", 9.0, 1.1, 4.1, 0.45, size=14, bold=True, color=ACCENT)
plan = [("0:00", "Dataset & RQs"),
        ("1:00", "Pipeline + Features"),
        ("2:00", "RQ1 — Supervised"),
        ("3:10", "RQ2 — Unsupervised"),
        ("4:20", "Insights & Ethics"),
        ("4:55", "Q&A")]
for i, (t, lbl) in enumerate(plan):
    y = 1.65 + i * 0.78
    rect(s, 9.0, y, 1.05, 0.55, ACCENT if i == 0 else DARK_BLUE)
    tb(s, t,   9.0,  y, 1.05, 0.52, size=12, bold=True,  color=DARK_BLUE if i==0 else ACCENT, align=PP_ALIGN.CENTER)
    tb(s, lbl, 10.1, y, 3.0,  0.52, size=12, color=WHITE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Dataset & Research Questions  [0:00 – 1:00]
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)
chrome(s, "0:00 – 1:00")
section_title(s, "Dataset & Research Questions")

# ── Left: dataset stats ───────────────────────────────────────────────────────
rect(s, 0.4, 1.55, 5.95, 5.25, LIGHT_GREY)
rect(s, 0.4, 1.55, 5.95, 0.44, DARK_BLUE)
tb(s, "Dataset at a Glance", 0.55, 1.58, 5.7, 0.38, size=13, bold=True, color=WHITE)

stats = [("Source",   "NYC TLC — Yellow Taxi Trip Records"),
         ("Period",   "June & July 2025"),
         ("Volume",   "~7.3 M trips · 1.4+ GB (Parquet)"),
         ("Extra",    "NYC Hourly Weather via Open-Meteo API"),
         ("Features", "14 clean features after wrangling"),
         ("Tools",    "Google Colab · pandas · scikit-learn · XGBoost")]
for i, (k, v) in enumerate(stats):
    y = 2.1 + i * 0.63
    tb(s, k+":", 0.55, y, 1.5,  0.5, size=12, bold=True,  color=MID_BLUE)
    tb(s, v,     2.05, y, 4.2,  0.5, size=12, color=DARK_GREY)

# ── Right: RQ cards ───────────────────────────────────────────────────────────
for ci, (rq, title, body, method) in enumerate([
    ("RQ1", "Supervised Classification",
     "Can we classify a high-tipper based on trip features, rate code & weather?",
     "Target: is_high_tip (tip ≥ 20 % of fare)  ·  Models: Random Forest, XGBoost"),
    ("RQ2", "Unsupervised Clustering",
     "Can we identify distinct natural clusters of trips based on temporal & financial features?",
     "Method: PCA  +  MiniBatch K-Means  (k = 7)"),
]):
    y = 1.55 + ci * 2.8
    rect(s, 6.75, y, 6.15, 2.55, LIGHT_BLUE)
    rect(s, 6.75, y, 6.15, 0.44, MID_BLUE)
    tb(s, rq + " — " + title, 6.9, y+0.04, 5.9, 0.38, size=13, bold=True, color=WHITE)
    tb(s, body,   6.9, y+0.56, 5.9, 0.9, size=12, italic=True, color=DARK_GREY)
    rect(s, 6.9, y+1.55, 5.95, 0.72, WHITE)
    tb(s, method, 7.0, y+1.6, 5.75, 0.58, size=12, bold=True, color=MID_BLUE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Pipeline + Feature Engineering  [1:00 – 2:00]
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)
chrome(s, "1:00 – 2:00")
section_title(s, "Pipeline Architecture & Feature Engineering")

# ── Pipeline flow (top half) ──────────────────────────────────────────────────
stages = [
    ("Data\nAcquisition",    MID_BLUE),
    ("Wrangling\n& EDA",     MID_BLUE),
    ("Feature\nEngineering", ACCENT),
    ("Supervised\nLearning", MID_BLUE),
    ("Unsupervised\nLearning",MID_BLUE),
    ("Insights\n& Report",   GREEN),
]
bw, bh, gap = 1.88, 1.35, 0.17
x0, y0 = 0.38, 1.58
for i, (lbl, col) in enumerate(stages):
    x = x0 + i*(bw+gap)
    rect(s, x, y0, bw, bh, LIGHT_GREY)
    rect(s, x, y0, bw, 0.42, col)
    tb(s, lbl, x+0.06, y0+0.03, bw-0.12, 0.36,
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(stages)-1:
        tb(s, "▶", x+bw+0.01, y0+0.45, gap+0.06, 0.44,
           size=16, color=MID_BLUE, align=PP_ALIGN.CENTER)

# Code-quality pills
pills = ["Modular functions", "Error handling", "Reproducible on new data", "Memory-efficient (dtype + GC)"]
for i, p in enumerate(pills):
    px2 = 0.38 + i*3.25
    rect(s, px2, y0+1.45, 3.1, 0.42, LIGHT_BLUE)
    rect(s, px2, y0+1.45, 0.07, 0.42, ACCENT)
    tb(s, p, px2+0.14, y0+1.48, 2.9, 0.36, size=11, color=DARK_BLUE)

# ── Feature engineering (bottom half) ────────────────────────────────────────
rect(s, 0.38, 3.35, 12.57, 0.04, LIGHT_BLUE)
tb(s, "Feature Engineering — 30 features → 12 FINAL_FEATURES (no leakage)",
   0.38, 3.45, 9.5, 0.42, size=13, bold=True, color=DARK_BLUE)

groups = [
    ("Transforms & Cyclical",
     "log_trip_distance · log_fare_amount\npickup_hour_sin/cos · pickup_day_sin/cos"),
    ("Domain Flags",
     "is_rush_hour · is_airport\nis_credit_card · fare_per_mile"),
    ("Polynomial / Interaction",
     "trip_distance²\ndistance×passengers · temp×precip"),
    ("Selection Funnel",
     "Filter (corr.) → Embedded (RF)\n→ Wrapper (RFE) → 12 features"),
]
for i, (gtitle, gtext) in enumerate(groups):
    gx = 0.38 + i*3.25
    rect(s, gx, 3.95, 3.1, 2.8, LIGHT_GREY)
    rect(s, gx, 3.95, 3.1, 0.4, MID_BLUE if i < 3 else GREEN)
    tb(s, gtitle, gx+0.1, 3.97, 2.9, 0.36, size=11, bold=True, color=WHITE)
    tb(s, gtext,  gx+0.1, 4.44, 2.9, 1.2,  size=11, color=DARK_GREY)

tb(s, "Excluded (leakage): tip_amount · total_amount · tip_ratio",
   0.38, 6.82, 9, 0.36, size=11, italic=True, color=RED)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — RQ1 Supervised Results  [2:00 – 3:10]
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)
chrome(s, "2:00 – 3:10")
section_title(s, "RQ1 — Supervised Classification Results")

# Metrics table
rect(s, 0.4, 1.55, 7.55, 0.42, DARK_BLUE)
for ci, (hdr, cx, cw) in enumerate([("Metric", 0.4, 3.1),
                                     ("Random Forest", 3.5, 2.2),
                                     ("XGBoost  ✓", 5.7, 2.25)]):
    tb(s, hdr, cx+0.05, 1.58, cw-0.1, 0.35,
       size=12, bold=True, color=WHITE,
       align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)

rows = [("CV / Val F1",    "0.8637 / 0.8636", "0.8637 / 0.8636"),
        ("Test Accuracy",  "0.77",             "0.77"),
        ("Recall High-Tip","1.00",             "1.00"),
        ("Recall Low-Tip", "0.08",             "0.08"),
        ("AUC-ROC",        "0.6293",           "0.6302"),
        ("Train Time",     "~250 s",           "~150 s  ✓")]
for ri, (m, rf, xg) in enumerate(rows):
    ry = 1.97 + ri*0.52
    bg = LIGHT_GREY if ri%2==0 else WHITE
    for ci, (val, cx, cw) in enumerate([(m,0.4,3.1),(rf,3.5,2.2),(xg,5.7,2.25)]):
        fc = RGBColor(0xEB,0xF5,0xFF) if ci==2 else bg
        rect(s, cx, ry, cw, 0.48, fc)
        tb(s, val, cx+0.07, ry+0.06, cw-0.14, 0.36,
           size=12, color=DARK_GREY if ci<2 else DARK_BLUE,
           bold=(ci==0), align=PP_ALIGN.LEFT if ci==0 else PP_ALIGN.CENTER)

rect(s, 0.4, 5.1, 7.55, 0.06, ACCENT)
rect(s, 0.4, 5.2, 7.55, 0.56, RGBColor(0xFF,0xF8,0xEC))
tb(s, "XGBoost selected — marginally higher AUC + 40 % faster training",
   0.55, 5.24, 7.2, 0.44, size=13, bold=True, color=ACCENT)

# Right insights panel
rect(s, 8.35, 1.55, 4.6, 4.2, LIGHT_BLUE)
rect(s, 8.35, 1.55, 4.6, 0.42, MID_BLUE)
tb(s, "Key Findings", 8.5, 1.58, 4.35, 0.36, size=13, bold=True, color=WHITE)
findings = [
    "Top features: fare_amount, improvement_surcharge, congestion_surcharge",
    "Tipping driven by payment structure — NOT weather or time of day",
    "Class imbalance 75/25 → Low-Tip recall only 0.08",
    "1 M stratified credit-card records · 70/15/15 split",
]
for fi, f in enumerate(findings):
    fy = 2.1 + fi*0.82
    rect(s, 8.5, fy, 0.07, 0.6, ACCENT)
    tb(s, f, 8.64, fy+0.04, 4.2, 0.55, size=11, color=DARK_GREY)

rect(s, 8.35, 5.6, 4.6, 0.88, RGBColor(0xEB,0xF5,0xFF))
tb(s, "Action: deploy as real-time dispatch scorer\nto flag high-tip trips for drivers.",
   8.5, 5.63, 4.35, 0.78, size=12, bold=True, color=DARK_BLUE)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — RQ2 Unsupervised Results  [3:10 – 4:20]
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)
chrome(s, "3:10 – 4:20")
section_title(s, "RQ2 — Unsupervised Clustering Results")

# Method bar
rect(s, 0.4, 1.55, 12.5, 0.5, LIGHT_BLUE)
tb(s, "PCA: 8 features → 6 PCs (≈90 % variance)   |   Elbow → k = 7   |   "
      "MiniBatch K-Means on ~7 M rows   |   Silhouette = 0.18   ·   DB = 1.35",
   0.55, 1.58, 12.1, 0.42, size=12, color=DARK_BLUE)

# Cluster table
col_w  = [0.65, 1.65, 4.15, 5.65]
col_x3 = [0.4,  1.05, 2.7,  6.85]
hdrs   = ["#", "Label", "Key Characteristic", "Actionable Strategy"]
rect(s, 0.4, 2.15, 12.5, 0.42, DARK_BLUE)
for hi, (hdr, cw, cx) in enumerate(zip(hdrs, col_w, col_x3)):
    tb(s, hdr, cx+0.05, 2.18, cw-0.1, 0.35,
       size=12, bold=True, color=WHITE)

clusters = [
    ("0",    "Long Trip",         "Highest mean distance — airport / inter-borough",  "Premium pricing · pre-position near JFK/LGA"),
    ("1–2",  "Short Trip",        "Low distance; likely artefact of k-choice",          "Investigate merging to sharpen short-trip targeting"),
    ("3",    "Medium Trip",       "Mid-range distance; no strong secondary feature",    "Standard base-fare offers"),
    ("4",    "Weather-Impacted",  "Very high mean precipitation",                       "Dynamic surge pricing during rain events"),
    ("5",    "Peak Morning",      "Lowest mean pickup hour — commuter segment",         "Pre-position in business districts before 08:00"),
    ("6",    "High-Occupancy",    "Highest mean passenger count — group trips",         "Group-ride promotions to raise revenue / vehicle"),
]
for ri, (cid, label, char, action) in enumerate(clusters):
    ry = 2.57 + ri*0.69
    rbg = WHITE if ri%2==0 else LIGHT_GREY
    for hi, (val, cw, cx) in enumerate(zip([cid,label,char,action], col_w, col_x3)):
        rect(s, cx, ry, cw, 0.65, rbg)
        fc = MID_BLUE if hi==1 else DARK_GREY
        tb(s, val, cx+0.07, ry+0.06, cw-0.14, 0.52,
           size=11, bold=(hi==1), color=fc)

tb(s, "Silhouette 0.18 is typical for high-dimensional real-world behavioural data — clusters are interpretable without tight separation.",
   0.4, 6.74, 12.5, 0.34, size=10, italic=True, color=DARK_GREY)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Insights, Ethics & Q&A  [4:20 – 5:00]
# ═════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.33, 7.5, WHITE)
chrome(s, "4:20 – 5:00")
section_title(s, "Insights, Ethics & Future Work")

# Four conclusion cards (left side)
cards = [
    (ACCENT,   "RQ1",       "XGBoost 77 % acc · AUC 0.63",       "Deploy real-time tip-scorer in dispatch app"),
    (MID_BLUE, "RQ2",       "7 clusters · clear archetypes",      "Surge (rain) · AM pre-positioning · group promos"),
    (RED,      "Limitation","Class imbalance 75/25 · AUC 0.63",   "Apply SMOTE + cost-sensitive training"),
    (GREEN,    "Scalability","1.4 GB in Colab via memory tricks",  "Migrate to cloud for monthly retraining"),
]
for i, (col, tag, finding, action) in enumerate(cards):
    y = 1.58 + i*1.3
    rect(s, 0.35, y, 0.12, 1.15, col)
    rect(s, 0.47, y, 6.8, 1.15, LIGHT_GREY if i%2==0 else WHITE)
    tb(s, tag,     0.65, y+0.06, 1.0,  0.38, size=13, bold=True,  color=col)
    tb(s, finding, 0.65, y+0.46, 5.85, 0.38, size=12, color=DARK_GREY)
    rect(s, 7.55, y+0.1, 3.8, 0.95, RGBColor(0xEB,0xF5,0xFF))
    tb(s, "▶ "+action, 7.65, y+0.13, 3.6, 0.78, size=11, bold=True, color=DARK_BLUE)

# Ethics column (far right)
rect(s, 11.55, 1.58, 1.4, 5.15, LIGHT_BLUE)
rect(s, 11.55, 1.58, 1.4, 0.42, MID_BLUE)
tb(s, "Ethics", 11.6, 1.61, 1.3, 0.36, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
ethics = ["Cash riders excluded", "Surge pricing fairness", "No PII stored", "Seasonal generalisability", "Clusters may shift monthly"]
for ei, e in enumerate(ethics):
    rect(s, 11.6, 2.1+ei*0.93, 0.07, 0.7, ACCENT)
    tb(s, e, 11.72, 2.12+ei*0.93, 1.15, 0.68, size=10, color=DARK_GREY)

# Q&A banner
rect(s, 0.35, 6.55, 11.12, 0.62, DARK_BLUE)
tb(s, "Thank you — Questions?", 0.55, 6.59, 10.7, 0.5,
   size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}  ({len(prs.slides)} slides)")
